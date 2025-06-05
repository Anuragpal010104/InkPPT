from pptx import Presentation

def split_bullets(summary, max_per_slide=4):
    # Split summary into bullet points (lines starting with * or -)
    lines = [line.strip() for line in summary.split('\n') if line.strip()]
    bullets = []
    for line in lines:
        if line.startswith('*') or line.startswith('-'):
            bullets.append(line)
        else:
            # If not a bullet, append to previous bullet (if exists)
            if bullets:
                bullets[-1] += ' ' + line
            else:
                bullets.append(line)
    # Group bullets into chunks
    return [bullets[i:i+max_per_slide] for i in range(0, len(bullets), max_per_slide)]

def create_presentation(structured_text, output_file="presentation.pptx"):
    prs = Presentation()
    for slide_data in structured_text:
        title = slide_data["title"]
        summary = slide_data["summary"]
        bullet_groups = split_bullets(summary)
        for idx, group in enumerate(bullet_groups):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            if idx == 0:
                slide.shapes.title.text = title
            else:
                # Remove the title placeholder for subsequent slides
                title_shape = slide.shapes.title
                if title_shape:
                    title_shape.text = ""
            content = slide.placeholders[1]
            content.text = '\n'.join(group)
    prs.save(output_file)